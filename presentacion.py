from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import date

# Información dinámica
ponente = "Presentacion 27000"
fecha_actual = date.today().strftime("%d/%m/%Y")

# Crear presentación
prs = Presentation()

# Tamaño de fuente
TITLE_FONT_SIZE = Pt(36)
CONTENT_FONT_SIZE = Pt(20)

# Función para añadir diapositiva
def add_slide(title, content_lines):
    slide_layout = prs.slide_layouts[1]  # Título y contenido
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    content_placeholder.text = '\n'.join(content_lines)

    # Cambiar tamaño del contenido
    for paragraph in content_placeholder.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = CONTENT_FONT_SIZE

# Diapositivas
slides_data = [
    ("ISO/IEC 27000: Entendiendo los Fundamentos de la Seguridad de la Información", [
        "Una Guía Esencial para la Implementación de un SGSI Robusto.",
        f"Ponente: {ponente}",
        f"Fecha: {fecha_actual}"
    ]),
    ("¿Por Qué la Seguridad de la Información es Crucial?", [
        "La Era Digital y sus Riesgos:",
        " - Dependencia creciente de la información digital.",
        " - Aumento de ciberataques (ransomware, phishing, robo de datos).",
        " - Regulaciones estrictas (GDPR, Ley de Protección de Datos Personales).",
        " - Impacto reputacional y financiero.",
        "¿Qué es la Seguridad de la Información?",
        " - Confidencialidad, Integridad, Disponibilidad (CID).",
        " - No solo tecnología: también personas, procesos y gestión.",
        "Rol de ISO/IEC 27000:",
        " - Marco común y vocabulario para gestionar la seguridad de la información."
    ]),
    ("La Familia de Normas ISO/IEC 27000", [
        "¿Qué es la Familia 27000?",
        " - Normas ISO/IEC para gestión de seguridad de la información.",
        " - Protege información sensible con enfoque sistemático.",
        "ISO/IEC 27000: Norma Introductoria (no certificable).",
        " - Visión general, vocabulario y principios.",
        "Objetivo: Ayudar a establecer, implementar y mejorar un SGSI."
    ]),
    ("ISO/IEC 27000: Alcance y Propósito", [
        "Título completo: \"Tecnología de la información - Técnicas de seguridad - SGSI - Visión general y vocabulario\"",
        "Propósito:",
        " - Descripción general de la serie 27000.",
        " - Definir terminología y conceptos fundamentales.",
        "Importancia:",
        " - Lenguaje común.",
        " - Facilita comprensión e implementación de SGSI.",
        " - Alinea prácticas de seguridad organizacionales."
    ]),
    ("Conceptos Clave de Seguridad de la Información", [
        " - Confidencialidad: Acceso solo a autorizados.",
        " - Integridad: Exactitud y completitud.",
        " - Disponibilidad: Acceso cuando se necesita.",
        " - Riesgo: Incertidumbre en los objetivos.",
        " - Activo: Elemento que requiere protección.",
        " - Amenaza: Causa de incidentes.",
        " - Vulnerabilidad: Debilidad explotable.",
        " - Control: Medida para modificar riesgo.",
        " - SGSI: Sistema de gestión de seguridad."
    ]),
    ("La Familia ISO/IEC 27000: Un Ecosistema de Normas", [
        " - 27000: Visión general y vocabulario.",
        " - 27001: Requisitos del SGSI (Certificable).",
        " - 27002: Código de prácticas.",
        " - 27003: Directrices de implementación.",
        " - 27004: Métricas y medición.",
        " - 27005: Gestión de riesgos.",
        " - 27007: Auditoría de SGSI.",
        " - 27017: Seguridad en la nube.",
        " - 27018: Protección de datos personales en la nube."
    ]),
    ("El Ciclo PDCA en un SGSI: Mejora Continua", [
        " - PLAN: Definir objetivos, procesos y recursos.",
        " - DO: Implementar y operar el SGSI.",
        " - CHECK: Monitorizar y auditar el desempeño.",
        " - ACT: Mejorar continuamente el SGSI."
    ]),
    ("Beneficios de un SGSI basado en ISO/IEC 27000", [
        " - Mejora la gestión de riesgos.",
        " - Cumplimiento normativo.",
        " - Protección reputacional.",
        " - Aumento de confianza.",
        " - Ventaja competitiva.",
        " - Reducción de costos.",
        " - Mejora continua.",
        " - Fomenta cultura de seguridad."
    ]),
    ("Rol Clave de ISO/IEC 27000 en la Implementación del SGSI", [
        " - Punto de referencia inicial.",
        " - Base para la capacitación.",
        " - Alineación con la dirección.",
        " - Guía para definir el alcance del SGSI."
    ]),
    ("Desafíos en la Implementación de un SGSI", [
        " - Falta de compromiso de la dirección.",
        " - Resistencia al cambio.",
        " - Falta de recursos.",
        " - Alcance mal definido.",
        " - Cultura organizacional.",
        " - Mantenimiento continuo.",
        " - Complejidad en la interpretación de requisitos."
    ]),
    ("Conclusión: La Seguridad como Pilar Estratégico", [
        " - ISO/IEC 27000 como base del SGSI.",
        " - No es solo certificación, es estrategia.",
        " - Compromiso continuo.",
        " - La inversión vale la pena."
    ]),
    ("Preguntas y Respuestas", [
        "¡Gracias por su atención!",
        "¿Alguna pregunta?"
    ])
]

# Agregar las diapositivas
for title, content in slides_data:
    add_slide(title, content)

# Guardar la presentación
prs.save("presentacion_ISO_27000.pptx")
print("✅ Presentación generada como 'presentacion_ISO_27000.pptx'")
